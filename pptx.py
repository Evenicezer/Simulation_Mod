from pptx import Presentation
from pptx.util import Inches

def draw_circle(presentation, radius):
    # Create a slide
    slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
    slide = presentation.slides.add_slide(slide_layout)

    # Calculate the center coordinates and diameter based on the radius
    center_x = Inches(5)  # Adjust the x-coordinate of the center as needed
    center_y = Inches(4)  # Adjust the y-coordinate of the center as needed
    diameter = radius * 2

    # Draw the circle shape
    shape = slide.shapes.add_shape(
        shape_type=MSO_SHAPE.OVAL,
        left=center_x - radius,
        top=center_y - radius,
        width=diameter,
        height=diameter
    )

    # Customize the circle shape
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Set the fill color (blue)
    shape.line.width = Pt(1)  # Set the outline width

# Example usage
radius = 3  # Replace with the desired radius based on the number of entries

# Create a PowerPoint presentation
presentation = Presentation()

# Draw a circle with the specified radius
draw_circle(presentation, radius)

# Save the presentation
presentation.save('circle.pptx')